#!/usr/bin/env -S uv run --script
# /// script
# requires-python = ">=3.12"
# dependencies = ["openpyxl", "python-pptx"]
# ///
"""Extract plain text from an .xlsx or .pptx file for ripgrep-all.

Reads the document from stdin (rga connects the file to stdin) and writes
extracted text to stdout. The first CLI argument selects the format.

Used as an rga custom adapter so `rga` can search inside spreadsheets and
slide decks, which its bundled pandoc adapter does not cover.
"""

import io
import sys


def extract_xlsx(data: bytes) -> str:
    """Return text from every cell of every sheet, one row per line."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"# Sheet: {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    return "\n".join(lines)


def extract_pptx(data: bytes) -> str:
    """Return text from every shape on every slide, tagged by slide number."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        lines.append(text)
    return "\n".join(lines)


def main() -> int:
    if len(sys.argv) < 2 or sys.argv[1] not in ("xlsx", "pptx"):
        sys.stderr.write("usage: office_extract.py {xlsx|pptx} < file\n")
        return 2
    data = sys.stdin.buffer.read()
    text = extract_xlsx(data) if sys.argv[1] == "xlsx" else extract_pptx(data)
    sys.stdout.write(text)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
